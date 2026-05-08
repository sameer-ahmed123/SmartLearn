########
#       SWITCHING FROM gemini-2.5-flash  to gemini-3-flash-preview
##
#       significantly faster response time for images
#       (down from 83-100 seconds to 13-16 seconds)

# NOTE gemini-3-flash-preview significantly fast ha gemini-2.5-flash se But
# Return tokens yah output significantly kam ha gemini-3-flash-preview ka
# COMPARE GENERATED LECTURE NUMBER 60 and 59
# 60 --- gemini-2.5-flash
# 59 --- gemini-3-flash-preview
##
########

import os
import uuid
import requests
import time
# type:ignore  its a part of the pymupdf library (required for pdf processing)
import fitz  # type:ignore
import logging  # type: ignore
from django.conf import settings # type:ignore
from PIL import Image  # type: ignore
from pptx import Presentation  # python-pptx library for PPTX # type: ignore
import google.genai as genai  # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE  # type:ignore
import cloudinary.uploader # type:ignore



logger = logging.getLogger(__name__)

api_key = settings.GEMINI_API_KEY
# Initialize the client, explicitly passing the key
if not api_key:
    raise ValueError(
        "GEMINI_API_KEY is missing in Django settings. Check your environment variables.")

# Initialize genai Client . Note the client will automatically read gemini api key from settings
client = genai.Client(api_key=api_key)

#######
##
# BELOW CODE : HELPER FUNCTIONS
##
######


def get_image_part(file_path):
    """ converts a local image file to a GenerativePart for Gemini API
        helps Ai model analize what the image contains
        Converts local image file to binary data (bytes) or phir bhje ga ai ko
    """
    with open(file_path, "rb") as f:
        image_bytes = f.read()
        # --- FIX MIME type ---
        # Use the correct standard MIME type for JPEG files.
        if file_path.lower().endswith(('.jpg', '.jpeg')):
            mime_type = "image/jpeg"
        elif file_path.lower().endswith('.png'):
            mime_type = "image/png"
        else:
            # Fallback or error handling for unsupported types
            raise ValueError(
                f"Unsupported image file extension for Gemini: {file_path}")
        return genai.types.Part.from_bytes(data=image_bytes, mime_type=mime_type)


def extract_content_from_pdf(full_path):
    """ 
    Extracts text content from a PDF using PyMuPDF.
    AND saves temporary images to feed to the ai model 
    returns text_content_str and [list_of_image_paths]
    """
    text_content = []
    image_paths = []

    # A temporary folder to store extracted images
    temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp_pdf_images')
    os.makedirs(temp_dir, exist_ok=True)

    try:
        with fitz.open(full_path) as doc:
            for page_num, page in enumerate(doc):
                # 1. extract text
                text_content.append(page.get_text())

                # 2. extract images
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]  # cross reference of image
                    base_image = doc.extract_image(xref)

                    if base_image:
                        image_bytes = base_image["image"]
                        # extention type ... jpg , png etc
                        ext = base_image["ext"]

                        # create unique filename
                        unique_id = uuid.uuid4().hex[:8]
                        image_filename = f"pdf_img_{page_num}_{img_index}_{unique_id}.{ext}"
                        temp_image_path = os.path.join(
                            temp_dir, image_filename)

                        # save the image bytes to temporary file
                        with open(temp_image_path, "wb") as img_file:
                            img_file.write(image_bytes)

                        image_paths.append(temp_image_path)

            extracted_text = "\n".join(text_content).strip()
        return extracted_text, image_paths
    except Exception as e:
        logger.error(f"Error extracting PDF Text From {full_path}: {e}")
        return "", []


def extract_content_from_pptx(full_path):
    """
    Extracts text and images  form a Power Point using python-pptx
    returns (text_content_str, [list_of_image_paths])
    """
    text_content = []
    image_paths = []

    temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp_pptx_images')
    os.makedirs(temp_dir, exist_ok=True)
    try:
        prs = Presentation(full_path)
        for slide_num, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                # 1. extract text from the presentation slides
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

                # 2. extract images from ppt (check for image shapes)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # checks if the encontered shape is even a picture or not
                    image_data = shape.image
                    image_bytes = image_data.blob
                    ext = image_data.ext  # ext means extenstion

                    # create a unique id for each image
                    unique_id = uuid.uuid4().hex[:8]
                    original_img_filename = f"pptx_img_{slide_num}_{shape_index}_{unique_id}.{ext}"
                    temp_image_path = os.path.join(
                        temp_dir, original_img_filename)

                    # save the image bytes to image_file
                    with open(temp_image_path, "wb") as img_file:
                        img_file.write(image_bytes)

                    # --- WMF/EMF CONVERSION STEP ---
                    # convert .wmb files to png (so ai can process it )
                    if ext.lower() in ['wmf', 'emf']:
                        final_image_path = temp_image_path
                        try:
                            png_filename = f"ppt_img_{slide_num}_{shape_index}_{unique_id}.png"
                            png_path = os.path.join(temp_dir, png_filename)
                            with Image.open(temp_image_path) as img:
                                # ensure its converted to RGB before storing in png
                                if img.mode not in ('RGB', 'RGBA'):
                                    img = img.convert('RGB')
                                img.save(png_path, 'PNG')
                            # use png images for LLM
                            final_image_path = png_path
                            # clean up unnessary files to save space
                            os.remove(temp_image_path)
                        except Exception as e:
                            logger.warning(
                                f"failed to convert WMF/EMF {temp_image_path} to png ,error: {e}")
                    else:
                        # If it's a standard format (PNG/JPG), use the path directly
                        final_image_path = temp_image_path

                    image_paths.append(final_image_path)
        extracted_text = "\n".join(text_content).strip()

        if not extracted_text and image_paths:
            logger.warning(
                f"PPTX extaction did not contain text or image for :{full_path}.")

        return extracted_text, image_paths
    except Exception as e:
        logger.error(f"Error extracting text from {full_path} : {e}")
        return "", []


def read_content_file(file_path_db):
    """
    1: Constructs the absolute path 
    2: determines the file type (e.g image ha yeh .txt), 
    3: read/processe the content.
    """
    if not file_path_db:
        return "Error: File path is None."

    full_path = os.path.join(settings.MEDIA_ROOT, file_path_db)

    if not os.path.exists(full_path):
        return f"Error: File not found at path {full_path}"

    # Use a try/except block to handle file reading errors
    try:
        if file_path_db.lower().endswith(('.png', '.jpg', '.jpeg')):
            return {
                "type": "image",
                "path": full_path,
                "content": f"load image file: {os.path.basename(file_path_db)}"
            }

        elif file_path_db.lower().endswith(('.txt', '.csv', '.md')):
            with open(full_path, 'r', encoding='utf-8') as f:
                return {'type': 'text', 'content': f.read()}

        elif file_path_db.lower().endswith('.pdf'):
            text_content, image_paths = extract_content_from_pdf(full_path)
            if not text_content and not image_paths:
                raise Exception("PDF extraction Failed")

            return {'type': 'multimodal', 'path': full_path, 'content': text_content, 'image_paths': image_paths}

        elif file_path_db.lower().endswith('.pptx'):
            text_content, image_paths = extract_content_from_pptx(full_path)
            if not text_content and not image_paths:
                raise Exception("PPTX extraction failed.")
            return {'type': 'multimodal', 'path': full_path, 'content': text_content, 'image_paths': image_paths}

        else:
            return {'type': 'unknown', 'content': f"Unsupported file type: {os.path.basename(file_path_db)}"}

    except Exception as e:
        return {'type': 'error', 'content': f"Error processing file {full_path}: {e}"}


#######
##
# BELOW CODE : BELONGS TO THE LECTURE GENERATION PIPELINE
##
######


def generate_lecture_script(prompt, file_data):
    """
    Calls Gemini API to Generate Script and Context in Json format
    """
    # 1. Prepare input Parts
    model_name = "gemini-2.5-flash"
    # model_name = "gemini-3-flash-preview"

    system_instructions = (
        "You are an expert curriculum designer. Your task is to generate a comprehensive"
        "lecture script and brief chatbot context based on the user's prompt and provided content."
        "The content may be text, images, or both. **If the text content is empty, rely on the "
        "attached images to generate the lecture.** "
        "the Output MUST be a JSON object with 'script' and 'context' keys."
    )
    # 2. build content array for the api call
    content_parts = [system_instructions, f"user prompt: {prompt}"]

    if file_data['type'] == 'image':
        # Assuming file_data['path'] holds the full path for the image
        # You will need to update read_content_file to return the path!
        image_part = get_image_part(file_data['path'])
        content_parts.append(image_part)
    elif file_data['type'] == 'text':
        content_parts.append(f"content for lecture: {file_data['content']}")
    elif file_data['type'] == 'multimodal':
        # new file type to send both iamge data and text content to ai (handles both pdf and ppt)

        # extracted text first
        content_parts.append(f"content for lecture: {file_data['content']}")

        # add all extracted (temporary) images as separate parts
        for image_path in file_data['image_paths']:
            try:
                # image_part helper use kar k image bytes or mime_type set kardo
                image_part = get_image_part(image_path)
                content_parts.append(image_part)
            except Exception as e:
                logger.warning(
                    f"failed to process extracted image {image_path}: {e}")
        content_parts.append(
            "IMPORTANT: Use provided text and ALL attached images to generate the script.")

    try:
        response = client.models.generate_content(
            model=model_name,
            contents=content_parts,
            config={
                "response_mime_type": "application/json",
                "response_schema": {
                    "type": "object",
                    "properties": {
                        "script": {"type": "string", "description": "The full, detailed Lecture Script."},
                        "context": {"type": "string", "description": "A short summary (max 200 words) of the material to train a Q&A chatbot"}
                    }
                }
            }
        )

        import json
        data = json.loads(response.text)
        return data.get("script"), data.get("context")
    except Exception as e:
        print(f"LLM API ERROR: {e}")
        return f"LLM Script Generation failed. {e}", "error context"



#######
##
# BELOW CODE : BELONGS TO THE ASSESMENT GENERATION PIPELINE
##
######


def generate_quiz_json(script, context, num_questions=5):
    """
    Generates a JSON object containing Multiple Choice Questions (MCQs) 
    based on the provided lecture script and context.
    """
    # 1. Setup Model
    # model_name = "gemini-2.5-flash"
    model_name = "gemini-3-flash-preview"

    system_instructions = (
        "You are an expert assessment creator. Your task is to generate a set of "
        f"{num_questions} multiple-choice questions based strictly on the provided "
        "Lecture Script and Context. \n"
        "- Questions should test conceptual understanding, not just rote memorization.\n"
        "- Provide exactly 4 options for each question.\n"
        "- Indicate the correct answer using the zero-based index (0, 1, 2, or 3)."
    )

    # 2. Prepare the Prompt Content
    user_prompt = (
        f"CONTEXT SUMMARY:\n{context}\n\n"
        f"FULL SCRIPT:\n{script}\n\n"
        f"TASK: Generate {num_questions} MCQs."
    )

    contents = [system_instructions, user_prompt]

    try:
        # 3. Call Gemini with specific Schema
        response = client.models.generate_content(
            model=model_name,
            contents=contents,
            config={
                "response_mime_type": "application/json",
                "response_schema": {
                    "type": "object",
                    "properties": {
                        "questions": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "question_text": {"type": "string"},
                                    "options": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Array of 4 option strings"
                                    },
                                    "correct_index": {
                                        "type": "integer",
                                        "description": "0-based index of the correct option"
                                    },
                                    "explanation": {
                                        "type": "string",
                                        "description": "Short explanation of why the answer is correct"
                                    }
                                },
                                "required": ["question_text", "options", "correct_index"]
                            }
                        }
                    }
                }
            }
        )

        import json
        data = json.loads(response.text)

        # Return just the list of questions
        return data.get("questions", [])

    except Exception as e:
        # log error in terminal
        logger.error(f"Quiz Generation LLM Error: {e}")
        # Re-raise the exception so the Celery task knows it failed
        raise Exception(f"Failed to generate quiz: {e}")


def generate_assignment_json(script, context,num_questions=5):
    """
    Generates a structured practical assignment based on the lecture script.
    """
    # model_name = "gemini-2.5-flash"
    model_name = "gemini-3-flash-preview"
    

    system_instructions = (
        "You are an expert university professor. Your task is to generate a practical, "
        "project-based assignment based strictly on the provided Lecture Script and Context.\n"
        "- The assignment should test application of knowledge, not just memorization.\n"
        f"- Define exactly {num_questions} actionable tasks for the student.\n"
        "- Provide a clear grading rubric mapped to these tasks.\n"
        "- Decide if the submission should be 'softcopy' (file upload/code) or 'hardcopy' (presentation/physical)."
    )

    user_prompt = (
        f"CONTEXT SUMMARY:\n{context}\n\n"
        f"FULL SCRIPT:\n{script}\n\n"
        f"TASK: Generate a practical assignment JSON."
    )

    contents = [system_instructions, user_prompt]

    try:
        response = client.models.generate_content(
            model=model_name,
            contents=contents,
            config={
                "response_mime_type": "application/json",
                "response_schema": {
                    "type": "object",
                    "properties": {
                        "title": {
                            "type": "string",
                            "description": "A professional title for the assignment"
                        },
                        "submission_type": {
                            "type": "string",
                            "description": "Either 'softcopy' or 'hardcopy'"
                        },
                        "tasks": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "List of 3-4 specific tasks the student must complete"
                        },
                        "rubric": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "criterion": {"type": "string", "description": "What is being graded"},
                                    "points": {"type": "integer", "description": "Points out of 100 total for this criterion"}
                                }
                            }
                        }
                    },
                    "required": ["title", "submission_type", "tasks", "rubric"]
                }
            }
        )

        import json
        data = json.loads(response.text)
        return data

    except Exception as e:
        logger.error(f"Assignment Generation LLM Error: {e}")
        raise Exception(f"Failed to generate assignment: {e}")    
    



def upload_video_to_cloudinary(video_source, lecture_id):
    """
    Uploads a video to Cloudinary and returns the secure URL and Public ID.
    Accepts both local file paths AND remote URLs (like HeyGen's temporary output link).
    """
    try:
        logger.info(f"Starting Cloudinary upload for Lecture {lecture_id}...")
        
        # use upload_large for videos because they can exceed standard file size limit
        # resource_type="video" is strictly required by Cloudinary for mp4 files
        response = cloudinary.uploader.upload_large(
            video_source,
            resource_type="video",
            folder="smartlearn_videos/lectures", 
            chunk_size=6000000 
        )
        
        secure_url = response.get('secure_url')
        public_id = response.get('public_id')
        
        logger.info(f"Cloudinary upload successful! URL: {secure_url}")
        return secure_url, public_id
        
    except Exception as e:
        logger.error(f"Cloudinary upload failed for Lecture {lecture_id}: {e}")
        return None, None
    


def generate_heygen_video(script_text):
    """
    1. Sends the script to HeyGen.
    2. Polls their server every 10 seconds until the video is done.
    3. Returns the temporary HeyGen download URL.
    """
    logger.info("Sending script to HeyGen API...")

    # --- PART 1: START GENERATION ---
    generate_url = "https://api.heygen.com/v2/video/generate"
    headers = {
        "X-Api-Key": settings.HEYGEN_API_KEY,
        "Content-Type": "application/json"
    }
    
    
    payload = {
        "video_inputs": [
            {
                "character": {
                    "type": "avatar",
                    "avatar_id": "Abigail_standing_office_front", 
                    "avatar_style": "normal"
                },
                "voice": {
                    "type": "text",
                    "input_text": script_text,
                    "voice_id": "1bd001e7e50f421d891986aad5158bc8" # Default female English voice
                },
            }
        ],
        "dimension": {
            "width": 1280,
            "height": 720
        }
    }
    
    response = requests.post(generate_url, json=payload, headers=headers)
    if response.status_code != 200:
        logger.error(f"HeyGen Start API Error: {response.text}")
        raise Exception(f"HeyGen Start API Error: {response.text}")
        
    data = response.json()
    video_id = data.get("data", {}).get("video_id")
    
    if not video_id:
        raise Exception("HeyGen did not return a video_id")

    logger.info(f"HeyGen generation started! Video ID: {video_id}. Waiting for completion...")

    # --- PART 2: POLL HEYGEN FOR VIDEO COMPLETION STATUS ---
    status_url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"
    headers_status = {"X-Api-Key": settings.HEYGEN_API_KEY}
    
    max_attempts = 45 # 45 attempts * 10 seconds = 7.5 minutes max wait time
    
    for attempt in range(max_attempts):
        status_res = requests.get(status_url, headers=headers_status)
        status_data = status_res.json()
        video_status = status_data.get("data", {}).get("status")
        
        if video_status == "completed":
            video_url = status_data.get("data", {}).get("video_url")
            logger.info("HeyGen rendering complete!")
            return video_url
            
        elif video_status == "failed":
            error_detail = status_data.get("data", {}).get("error", "Unknown error")
            raise Exception(f"HeyGen rendering failed internally: {error_detail}")
        
        logger.info(f"HeyGen status: {video_status}... checking again in 10 seconds.")
        time.sleep(10) # Pause the Celery worker for 10 seconds before asking again
        
    raise Exception("HeyGen video generation timed out (exceeded 7.5 minutes).")


def generate_did_video(script_text):
    """
    1. Sends a static image and script to D-ID.
    2. Polls their server until the video is done.
    3. Returns the temporary D-ID download URL.
    """
    logger.info("Sending script to D-ID API...")

    # --- PART 1: START GENERATION ---
    generate_url = "https://api.d-id.com/talks"
    
    # D-ID requires the key to be passed as Basic Auth. 
    # When you copy your key from the D-ID dashboard, it will likely be in the correct format.
    headers = {
        "accept": "application/json",
        "content-type": "application/json",
        "authorization": f"Basic {settings.DID_API_KEY}"
    }
    
    payload = {
        "script": {
            "type": "text",
            "input": script_text,
            "provider": {
                "type": "microsoft",
                "voice_id": "en-US-JennyNeural" # Standard professional female voice
            }
        },
        # can replace this with any public URL of a clean, forward-facing headshot!
        # the avatar source image
        "source_url": "https://plus.unsplash.com/premium_photo-1661505218403-c684557a824d?q=80&w=687&auto=format&fit=crop&ixlib=rb-4.1.0&ixid=M3wxMjA3fDB8MHxwaG90by1wYWdlfHx8fGVufDB8fHx8fA%3D%3D" 
    }
    
    response = requests.post(generate_url, json=payload, headers=headers)
    
    if response.status_code != 201:
        logger.error(f"D-ID Start API Error: {response.text}")
        raise Exception(f"D-ID Start API Error: {response.text}")
        
    data = response.json()
    talk_id = data.get("id")
    
    if not talk_id:
        raise Exception("D-ID did not return a talk_id")

    logger.info(f"D-ID generation started! Talk ID: {talk_id}. Waiting for completion...")

    # --- PART 2: POLL D-ID FOR VIDEO COMPLETION STATUS  ---
    status_url = f"https://api.d-id.com/talks/{talk_id}"
    
    max_attempts = 45 # 45 attempts * 10 seconds = 7.5 minutes max wait time
    
    for attempt in range(max_attempts):
        status_res = requests.get(status_url, headers=headers)
        status_data = status_res.json()
        video_status = status_data.get("status")
        
        if video_status == "done":
            video_url = status_data.get("result_url")
            logger.info("D-ID rendering complete!")
            return video_url
            
        elif video_status == "error":
            error_detail = status_data.get("last_error", "Unknown error")
            raise Exception(f"D-ID rendering failed internally: {error_detail}")
        
        logger.info(f"D-ID status: {video_status}... checking again in 10 seconds.")
        time.sleep(10) 
        
    raise Exception("D-ID video generation timed out (exceeded 7.5 minutes).")
